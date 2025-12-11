#%%

from pptx import Presentation
from pptx.util import Inches
import os
import glob

#%%

# Directory where your images are stored
actual_data_base_dir = "/shared/ad150/event3d/"
# image_dir = "/shared/ad150/event3d/marigold/checkpoint/train_marigold_monocular/visualization/iter_007250/"
image_dir = "/shared/ad150/event3d/marigold/checkpoint/train_marigold_monocular/visualization/iter_008000/"
output_dir = "./"

image_files = []
image_files_magma = []
# os.walk returns the [current folder path, the subfolders in it, the files in it]
# It will "walk" through the base directory so that the all the "current folder path"s will cover
# all possible subdirectories
for root, dirs, files in os.walk(image_dir):
    for file in files:
        if file.endswith(('.png', '.jpg', '.jpeg')):
            if ("magma" in file):
                image_files_magma.append(os.path.relpath(os.path.join(root, file), image_dir))
            else:
                image_files.append(os.path.relpath(os.path.join(root, file), image_dir))

image_files.sort()
image_files_magma.sort()

print(image_files)
print(image_files_magma)

#%%

# the rgb file paths can be retrieved from the image_files list

rgb_files = []
for file_iter in range(len(image_files)):
    print(file_iter)
    if ("carla" in image_files[file_iter]):
        rgb_path = (image_files[file_iter].replace("_", "/").replace(".tif", "")
                    .replace("Camera/1", "Camera_1").replace("Camera/0", "Camera_0")
                    .replace("event/LINEAR/event/frame/", "rgb_").replace("/magma", "")
                    .replace("/vis", "").replace("vkitti", "vkitti2")
                    .replace("sequence/", "sequence_")
                    .replace("/val", "_val").replace("/train", "_train")
                    .replace("events/frames", "rgb/data"))
        # This will have the correct directory structure for the rgb files
        # but not the actual file name
        # find the correct file name by checking the actual directory with correct number in it
        
        # get rid of the last part of the path
        rgb_path_dir = rgb_path[:rgb_path.rfind("/")]
        if False == os.path.exists(os.path.join(actual_data_base_dir, rgb_path_dir)):
            rgb_path_dir = rgb_path_dir.replace("data", "frames")

        filename = rgb_path[rgb_path.rfind("/")+1:]
        # get only the number part of the filename
        number = int(filename[filename.rfind("_")+1:filename.rfind(".")])

        # now search for the correct file name which has the same number
        filenames = glob.glob(os.path.join(actual_data_base_dir, rgb_path_dir, f"*{number:04d}*"))
        try:
            rgb_path = filenames[0]
        except:
            print(f"-- {image_files[file_iter]}--")
            print(f"-- {rgb_path_dir}--")
            print(f"-- {number}--")
    elif ("vkitti" in image_files[file_iter]):
        rgb_path = (image_files[file_iter].replace("_", "/").replace(".tif", "")
                    .replace("Camera/1", "Camera_1").replace("Camera/0", "Camera_0")
                    .replace("event/LINEAR/event/frame/", "rgb_").replace("/magma", "")
                    .replace("/vis", "").replace("vkitti", "vkitti2")
                    .replace("png", "jpg"))

    rgb_path = os.path.join(actual_data_base_dir, rgb_path)

    print(image_files[file_iter])
    print(rgb_path)
    # check if the rgb file exists
    rgb_exists = os.path.exists(rgb_path)
    print(rgb_exists)

    rgb_files.append(rgb_path)

#%%

# the depth file paths can be retrieved from the image_files list

depth_files = []
for file_iter in range(len(image_files)):
    print(file_iter)
    if ("carla" in image_files[file_iter]):
        depth_path = (image_files[file_iter].replace("_", "/").replace(".tif", "")
                    .replace("Camera/1", "Camera_1").replace("Camera/0", "Camera_0")
                    .replace("event/LINEAR/event/frame/", "depth_").replace("/magma", "")
                    .replace("/vis", "").replace("vkitti", "vkitti2")
                    .replace("sequence/", "sequence_")
                    .replace("/val", "_val").replace("/train", "_train")
                    .replace("events/frames", "depth/frames"))
        # This will have the correct directory structure for the rgb files
        # but not the actual file name
        # find the correct file name by checking the actual directory with correct number in it
        
        # get rid of the last part of the path
        depth_path_dir = depth_path[:depth_path.rfind("/")]
        if False == os.path.exists(os.path.join(actual_data_base_dir, depth_path_dir)):
            depth_path_dir = depth_path_dir.replace("data", "frames")

        filename = depth_path[depth_path.rfind("/")+1:]
        # get only the number part of the filename
        number = int(filename[filename.rfind("_")+1:filename.rfind(".")])

        # now search for the correct file name which has the same number
        extension = ".png"
        filenames = glob.glob(os.path.join(actual_data_base_dir, depth_path_dir, f"*{number:04d}*{extension}"))
        try:
            depth_path = filenames[0]
        except:
            print(f"-- {image_files[file_iter]}--")
            print(f"-- {depth_path_dir}--")
            print(f"-- {number}--")
    elif ("vkitti" in image_files[file_iter]):
        depth_path = (image_files[file_iter].replace("_", "/").replace(".tif", "")
                    .replace("Camera/1", "Camera_1").replace("Camera/0", "Camera_0")
                    .replace("event/LINEAR/event/frame/", "depth_").replace("/magma", "")
                    .replace("/vis", "").replace("vkitti", "vkitti2")
                    .replace("rgb", "depth")
                    .replace("jpg", "png"))

    depth_path = os.path.join(actual_data_base_dir, depth_path)

    print(image_files[file_iter])
    print(depth_path)
    # check if the rgb file exists
    depth_exists = os.path.exists(depth_path)
    print(depth_exists)

    depth_files.append(depth_path)

#%%

# the event file paths can be retrieved from the image_files list

event_files = []
for file_iter in range(len(image_files)):
    print(file_iter)
    if ("carla" in image_files[file_iter]):
        event_path = (image_files[file_iter].replace("_", "/").replace(".tif", "")
                    .replace("Camera/1", "Camera_1").replace("Camera/0", "Camera_0")
                    .replace("/magma", "")
                    .replace("/vis", "").replace("vkitti", "vkitti2")
                    .replace("sequence/", "sequence_")
                    .replace("/val", "_val").replace("/train", "_train")
                    .replace("event/frame/", "event_frame_")
                    .replace("frames/event/", "frames_event/")
                    .replace("png", "tif"))
        # This will have the correct directory structure for the rgb files
        # but not the actual file name
        # find the correct file name by checking the actual directory with correct number in it
        
        # get rid of the last part of the path
        event_path_dir = event_path[:event_path.rfind("/")]
        if False == os.path.exists(os.path.join(actual_data_base_dir, event_path_dir)):
            event_path_dir = event_path_dir.replace("data", "frames")

        filename = event_path[event_path.rfind("/")+1:]
        # get only the number part of the filename
        number = int(filename[filename.rfind("_")+1:filename.rfind(".")])

        # now search for the correct file name which has the same number
        filenames = glob.glob(os.path.join(actual_data_base_dir, event_path_dir, f"*{number:04d}*"))
        try:
            event_path = filenames[0]
        except:
            print(f"-- {image_files[file_iter]}--")
            print(f"-- {event_path_dir}--")
            print(f"-- {number}--")
    elif ("vkitti" in image_files[file_iter]):
        event_path = (image_files[file_iter].replace("_", "/").replace(".tif", "")
                    .replace("Camera/1/event", "Camera_1_event")
                    .replace("Camera/0/event", "Camera_0_event")
                    .replace("/magma", "")
                    .replace("/vis", "").replace("vkitti", "vkitti2")
                    .replace("event/frame/", "event_frame_")
                    .replace("png", "tif")
                    .replace("jpg", "tif"))

    event_path = os.path.join(actual_data_base_dir, event_path)

    print(image_files[file_iter])
    print(event_path)
    # check if the rgb file exists
    event_exists = os.path.exists(event_path)
    print(event_exists)

    event_files.append(event_path)


#%%
# Create a new PowerPoint presentation
prs = Presentation()  # Create a new PowerPoint presentation
prs.slide_height = Inches(9)
prs.slide_width = Inches(18)

inch = Inches(1)
slide_height = prs.slide_height/inch
slide_width = prs.slide_width/inch
print(slide_height, slide_width)

# Loop over the images and add a set of (image, image_magma, rgb) to a new slide
for file_iter in range(len(image_files)):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
    image_path = os.path.join(image_dir, image_files[file_iter])
    image_path_magma = os.path.join(image_dir, image_files_magma[file_iter])
    rgb_path = rgb_files[file_iter]
    depth_path = depth_files[file_iter]
    event_path = event_files[file_iter]

    # Add 5 images - 2 for the depth, 
    # 1 for the rgb, 1 for GT depth, 1 for the input event image
    slide.shapes.add_picture(depth_path, Inches(0.5), Inches(0.2), width=Inches(5))  # Position and resize
    slide.shapes.add_picture(image_path, Inches(6.5), Inches(0.2), width=Inches(5))
    slide.shapes.add_picture(image_path_magma, Inches(12.5), Inches(0.2), width=Inches(5))
    slide.shapes.add_picture(rgb_path, Inches(3.5), Inches(4.5), width=Inches(5))
    slide.shapes.add_picture(event_path, Inches(9.5), Inches(4.5), width=Inches(5))

# Loop over the images and add each one to a new slide
# for image_file in image_files:
#     slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
#     image_path = os.path.join(image_dir, image_file)
#     slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(5))

# Save the presentation
vis_num = int(image_dir[image_dir.rfind("_")+1:image_dir.rfind("/")])
savename = os.path.join(output_dir, f"results_presentation_vis_{vis_num}.pptx")
prs.save(savename)