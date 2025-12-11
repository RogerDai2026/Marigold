#%%

from pptx import Presentation
from pptx.util import Inches
import os
import glob

import numpy as np
import imageio.v3 as io
import matplotlib.pyplot as plt

import datetime
import argparse


if "__main__" == __name__:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--actual_data_base_dir",
        type=str,
        default="/shared/ad150/event3d/",
    )
    parser.add_argument(
        "--image_dir",
        type=str,
        required=True,
        help="Directory where the images are stored, example: "
        "/shared/ad150/event3d/marigold/checkpoint/train_marigold_monocular/visualization/iter_010250/",
    )
    parser.add_argument(
        "--output_dir",
        type=str,
        default="/shared/ad150/event3d/",
    )    


    args = parser.parse_args()

    actual_data_base_dir = args.actual_data_base_dir
    image_dir = args.image_dir
    output_dir = args.output_dir

    # if the image_dir does not end with a slash, add it
    if not image_dir.endswith("/"):
        image_dir += "/"

    # # Directory where your images are stored
    # actual_data_base_dir = "/shared/ad150/event3d/"
    # image_dir = "/shared/ad150/event3d/marigold/checkpoint/train_marigold_monocular/visualization/iter_010250/"
    # output_dir = "/shared/ad150/event3d/"

    exp_name = image_dir[image_dir.find("checkpoint/")+11:]
    exp_name = exp_name[:exp_name.find("/")]
    print(f"Experiment name: {exp_name}")

    image_files = []
    image_files_magma = []
    # os.walk returns the [current folder path, the subfolders in it, the files in it]
    # It will "walk" through the base directory so that the all the "current folder path"s will cover
    # all possible subdirectories
    for root, dirs, files in os.walk(image_dir):
        for file in files:
            if file.endswith(('.png', '.jpg', '.jpeg')):
                if ("magma" in file):
                    image_files_magma.append(os.path.relpath(os.path.join(root, file), image_dir))
                else:
                    image_files.append(os.path.relpath(os.path.join(root, file), image_dir))

    image_files.sort()
    image_files_magma.sort()

    #print(image_files)
    #print(image_files_magma)

    #%%

    # the rgb file paths can be retrieved from the image_files list

    rgb_files = []
    for file_iter in range(len(image_files)):
        #print(file_iter)
        if ("carla" in image_files[file_iter]):
            rgb_path = (image_files[file_iter].replace("_", "/").replace(".tif", "")
                        .replace("Camera/1", "Camera_1").replace("Camera/0", "Camera_0")
                        .replace("event/LINEAR/event/frame/", "rgb_").replace("/magma", "")
                        .replace("/vis", "").replace("vkitti", "vkitti2")
                        .replace("sequence/", "sequence_")
                        .replace("/val", "_val").replace("/train", "_train")
                        .replace("events/frames", "rgb/data"))
            # This will have the correct directory structure for the rgb files
            # but not the actual file name
            # find the correct file name by checking the actual directory with correct number in it
            
            # get rid of the last part of the path
            rgb_path_dir = rgb_path[:rgb_path.rfind("/")]
            if False == os.path.exists(os.path.join(actual_data_base_dir, rgb_path_dir)):
                rgb_path_dir = rgb_path_dir.replace("data", "frames")

            filename = rgb_path[rgb_path.rfind("/")+1:]
            # get only the number part of the filename
            number = int(filename[filename.rfind("_")+1:filename.rfind(".")])

            # now search for the correct file name which has the same number
            filenames = glob.glob(os.path.join(actual_data_base_dir, rgb_path_dir, f"*{number:04d}*"))
            try:
                rgb_path = filenames[0]
            except:
                print(f"-- {image_files[file_iter]}--")
                print(f"-- {rgb_path_dir}--")
                print(f"-- {number}--")
        elif ("vkitti" in image_files[file_iter]):
            rgb_path = (image_files[file_iter].replace("_", "/").replace(".tif", "")
                        .replace("Camera/1", "Camera_1").replace("Camera/0", "Camera_0")
                        .replace("event/LINEAR/event/frame/", "rgb_").replace("/magma", "")
                        .replace("/vis", "").replace("vkitti", "vkitti2")
                        .replace("png", "jpg"))

        rgb_path = os.path.join(actual_data_base_dir, rgb_path)

        #print(image_files[file_iter])
        #print(rgb_path)
        # check if the rgb file exists
        rgb_exists = os.path.exists(rgb_path)
        #print(rgb_exists)

        rgb_files.append(rgb_path)

    #%%

    # the depth file paths can be retrieved from the image_files list

    depth_files = []
    for file_iter in range(len(image_files)):
        #print(file_iter)
        if ("carla" in image_files[file_iter]):
            depth_path = (image_files[file_iter].replace("_", "/").replace(".tif", "")
                        .replace("Camera/1", "Camera_1").replace("Camera/0", "Camera_0")
                        .replace("event/LINEAR/event/frame/", "depth_").replace("/magma", "")
                        .replace("/vis", "").replace("vkitti", "vkitti2")
                        .replace("sequence/", "sequence_")
                        .replace("/val", "_val").replace("/train", "_train")
                        .replace("events/frames", "depth/frames"))
            # This will have the correct directory structure for the rgb files
            # but not the actual file name
            # find the correct file name by checking the actual directory with correct number in it
            
            # get rid of the last part of the path
            depth_path_dir = depth_path[:depth_path.rfind("/")]
            if False == os.path.exists(os.path.join(actual_data_base_dir, depth_path_dir)):
                depth_path_dir = depth_path_dir.replace("data", "frames")

            filename = depth_path[depth_path.rfind("/")+1:]
            # get only the number part of the filename
            number = int(filename[filename.rfind("_")+1:filename.rfind(".")])

            # now search for the correct file name which has the same number
            extension = ".png"
            filenames = glob.glob(os.path.join(actual_data_base_dir, depth_path_dir, f"*{number:04d}*{extension}"))
            try:
                depth_path = filenames[0]
            except:
                print(f"-- {image_files[file_iter]}--")
                print(f"-- {depth_path_dir}--")
                print(f"-- {number}--")
        elif ("vkitti" in image_files[file_iter]):
            depth_path = (image_files[file_iter].replace("_", "/").replace(".tif", "")
                        .replace("Camera/1", "Camera_1").replace("Camera/0", "Camera_0")
                        .replace("event/LINEAR/event/frame/", "depth_").replace("/magma", "")
                        .replace("/vis", "").replace("vkitti", "vkitti2")
                        .replace("rgb", "depth")
                        .replace("jpg", "png"))

        depth_path = os.path.join(actual_data_base_dir, depth_path)

        #print(image_files[file_iter])
        #print(depth_path)
        # check if the rgb file exists
        depth_exists = os.path.exists(depth_path)
        #print(depth_exists)

        depth_files.append(depth_path)

    #%%

    # the event file paths can be retrieved from the image_files list

    event_files = []
    for file_iter in range(len(image_files)):
        #print(file_iter)
        if ("carla" in image_files[file_iter]):
            event_path = (image_files[file_iter].replace("_", "/").replace(".tif", "")
                        .replace("Camera/1", "Camera_1").replace("Camera/0", "Camera_0")
                        .replace("/magma", "")
                        .replace("/vis", "").replace("vkitti", "vkitti2")
                        .replace("sequence/", "sequence_")
                        .replace("/val", "_val").replace("/train", "_train")
                        .replace("event/frame/", "event_frame_")
                        .replace("frames/event/", "frames_event/")
                        .replace("png", "tif"))
            # This will have the correct directory structure for the rgb files
            # but not the actual file name
            # find the correct file name by checking the actual directory with correct number in it
            
            # get rid of the last part of the path
            event_path_dir = event_path[:event_path.rfind("/")]
            if False == os.path.exists(os.path.join(actual_data_base_dir, event_path_dir)):
                event_path_dir = event_path_dir.replace("data", "frames")

            filename = event_path[event_path.rfind("/")+1:]
            # get only the number part of the filename
            number = int(filename[filename.rfind("_")+1:filename.rfind(".")])

            # now search for the correct file name which has the same number
            filenames = glob.glob(os.path.join(actual_data_base_dir, event_path_dir, f"*{number:04d}*"))
            try:
                event_path = filenames[0]
            except:
                print(f"-- {image_files[file_iter]}--")
                print(f"-- {event_path_dir}--")
                print(f"-- {number}--")
        elif ("vkitti" in image_files[file_iter]):
            event_path = (image_files[file_iter].replace("_", "/").replace(".tif", "")
                        .replace("Camera/1/event", "Camera_1_event")
                        .replace("Camera/0/event", "Camera_0_event")
                        .replace("/magma", "")
                        .replace("/vis", "").replace("vkitti", "vkitti2")
                        .replace("event/frame/", "event_frame_")
                        .replace("png", "tif")
                        .replace("jpg", "tif"))

        event_path = os.path.join(actual_data_base_dir, event_path)

        #print(image_files[file_iter])
        #print(event_path)
        # check if the rgb file exists
        event_exists = os.path.exists(event_path)
        #print(event_exists)

        event_files.append(event_path)


    #%%

    # get current timestamp and use that in the file name
    # get time now
    now = datetime.datetime.now()
    timestamp = now.strftime("%Y%m%d_%H%M%S")
    vis_num = int(image_dir[image_dir.rfind("_")+1:image_dir.rfind("/")])

    # Create a new PowerPoint presentation
    prs = Presentation()  # Create a new PowerPoint presentation
    prs.slide_height = Inches(9)
    prs.slide_width = Inches(18)

    inch = Inches(1)
    slide_height = prs.slide_height/inch
    slide_width = prs.slide_width/inch
    #print(slide_height, slide_width)

    def convert_to_suitable_format(img, h, w):
        if (img.shape[0] != h) or (img.shape[1] != w):
            img = img[:h, :w]
        if (len(img.shape) == 2):
            img = img[:,:,None].repeat(3, axis=-1)
        if img.dtype == np.float32:
            img = (img * 255).astype(np.uint8)
        if img.dtype == np.uint16:
            img = (img/256).astype(np.uint8)
        return img

    def convert_to_inverse_log_magma(img, h, w, amax, amin=None):
        img_temp = img[:h, :w]
        img_temp = np.clip(img_temp, a_min = amin, a_max=amax)
        img_temp = np.log(1/img_temp)
        img_temp = (img_temp - img_temp.min())/np.abs(img_temp.max() - img_temp.min())
        img_temp = plt.cm.magma(img_temp)[:,:,:3]
        img_magma = (img_temp * 255).astype(np.uint8)
        return img_magma

    # Loop over the images and add a set of (image, image_magma, rgb) to a new slide
    for file_iter in range(len(image_files)):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
        image_path = os.path.join(image_dir, image_files[file_iter])
        image_path_magma = os.path.join(image_dir, image_files_magma[file_iter])
        rgb_path = rgb_files[file_iter]
        depth_path = depth_files[file_iter]
        event_path = event_files[file_iter]

        # Add 5 images - 2 for the depth, 
        # 1 for the rgb, 1 for GT depth, 1 for the input event image
        # slide.shapes.add_picture(depth_path, Inches(0.5), Inches(0.2), width=Inches(5))  # Position and resize
        # slide.shapes.add_picture(image_path, Inches(6.5), Inches(0.2), width=Inches(5))
        # slide.shapes.add_picture(image_path_magma, Inches(12.5), Inches(0.2), width=Inches(5))
        # slide.shapes.add_picture(rgb_path, Inches(3.5), Inches(4.5), width=Inches(5))
        # slide.shapes.add_picture(event_path, Inches(9.5), Inches(4.5), width=Inches(5))

        # Also make a single image for all the images combined
        img_image = io.imread(image_path)
        h = img_image.shape[0]
        w = img_image.shape[1]
        img_image_magma = io.imread(image_path_magma)
        img_rgb = io.imread(rgb_path)
        img_event = io.imread(event_path)
        img_depth = io.imread(depth_path)
        # Also render a magma version of the depth image
        # if dataset is vkitti, amax is 25000, if carla, amax is 250
        if "vkitti" in depth_path:
            amax = 25000
            amin = 5
            factor = 256
        else:
            amax = 25000
            amin = 5
            factor = 1

        # Convert images into suitable format
        depth_GT = convert_to_suitable_format(img_depth, h, w)
        depth_pred = convert_to_suitable_format(img_image, h, w)

        img_depth_magma = convert_to_inverse_log_magma(img_depth/factor, h, w, amin=amin, amax=amax)
        img_image_magma = convert_to_inverse_log_magma(img_image/256, h, w, amin=amin, amax=amax)


        # Combine image in a single 2x3 grid
        img_combined = np.zeros((img_image.shape[0]*2, img_image.shape[1]*3, 3), dtype=np.uint8)
        img_combined[:img_image.shape[0], :img_image.shape[1],:] = depth_GT
        img_combined[:img_image.shape[0], img_image.shape[1]:2*img_image.shape[1]] = img_depth_magma#img_depth_magma[:img_image.shape[0], :img_image.shape[1],:]
        img_combined[:img_image.shape[0], 2*img_image.shape[1]:] = img_event[:img_image.shape[0], :img_image.shape[1],:]
        img_combined[img_image.shape[0]:, :img_image.shape[1]] = depth_pred
        img_combined[img_image.shape[0]:, img_image.shape[1]:2*img_image.shape[1]] = img_image_magma
        img_combined[img_image.shape[0]:, 2*img_image.shape[1]:] = img_rgb[:img_image.shape[0], :img_image.shape[1],:]

        out_new = "view_results_here/"
        a = image_files[file_iter]
        number = a[a.rfind("_")+1:a.find(".")]


        save_dir = os.path.join(output_dir, out_new, exp_name, f"{vis_num}")
        if (not os.path.exists(save_dir)):
            os.makedirs(save_dir)

        plt.figure(figsize=(16, 9))
        plt.imshow(img_combined)
        plt.axis('off')
        plt.tight_layout()
        plt.title(image_files[file_iter])
        # Add the combined image to the slide
        plt.savefig(f"{save_dir}/temp_{number}.tif", bbox_inches='tight', pad_inches=0)
        slide.shapes.add_picture(f"{save_dir}/temp_{number}.tif", Inches(0.5), Inches(0.5), width=Inches(16))
        



    # Save the presentation
    savename = os.path.join(output_dir, f"{save_dir}/results_presentation_vis_{vis_num}.pptx")
    prs.save(savename)

    print(f"Saved the images and presentation to {save_dir}")